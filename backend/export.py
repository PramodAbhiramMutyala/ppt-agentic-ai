import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def apply_theme_styling(slide, title_text, is_title_slide=False):
    """
    Applies custom geometric themes and colors since python-pptx uses blank slides by default.
    """
    # Create a nice header banner block for content slides
    if not is_title_slide:
        banner = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0), Inches(10), Inches(1.2)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(34, 45, 65) # Dark slate blue
        banner.line.color.rgb = RGBColor(34, 45, 65)
        
        # Add custom title text block overlayed on the banner
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(36)
        p.font.bold = True
    else:
        # Style Title Slide
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(34, 45, 65)
        
def create_presentation(topic: str, slide_data: list[dict], filename: str = None) -> str:
    prs = Presentation()
    # Widescreen 16:9 ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 1. Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_slide_layout)
    apply_theme_styling(title_slide, "", is_title_slide=True)
    
    txBox = title_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{topic.title()}"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(54)
    p.font.bold = True
    
    subtitle = title_slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    p_sub = subtitle.text_frame.paragraphs[0]
    p_sub.text = "Generated autonomously by AI Agents"
    p_sub.font.color.rgb = RGBColor(200, 200, 200)
    p_sub.font.size = Pt(24)
    
    # 2. Add Content Slides
    for item in slide_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        title_str = item.get("title", "Untitled Slide")
        apply_theme_styling(slide, title_str, is_title_slide=False)
        
        # Add Bullets block on the Left
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for idx, bullet in enumerate(item.get("bullets", [])):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.space_after = Pt(14)
            p.font.color.rgb = RGBColor(50, 50, 50)
            
        # Add the Auto-Generated Image on the Right
        img_path = item.get("image_path")
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(5.2), Inches(1.4), width=Inches(4.3))
            except Exception as e:
                print(f"Failed to attach image to slide: {e}")
                
        # Optional: Add image query text at the bottom as speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = f"Suggested Image Generation Prompt:\n{item.get('image_query', '')}"
        
    # Save Presentation
    if not filename:
        clean_topic = "".join(x for x in topic if x.isalnum() or x in " -_").strip().replace(" ", "_")
        filename = f"{clean_topic}_Presentation.pptx"
        
    prs.save(filename)
    return filename
